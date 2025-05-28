import os
import re
from pptx import Presentation
from docx import Document


def clean_xml_text(text):
    """Usuwa niedozwolone znaki XML (null, kontrolne)"""
    return re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F]', '', text or '')


def extract_notes_from_pptx(pptx_path):
    """Zwraca listę (Slide X, notatka) z pliku PPTX, pomijając ukryte slajdy"""
    prs = Presentation(pptx_path)
    notes = []

    for idx, slide in enumerate(prs.slides):
        slide_number = idx + 1

        # Sprawdź, czy slajd jest ukryty
        is_hidden = slide._element.get("show") == "0"
        if is_hidden:
            continue

        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            note = slide.notes_slide.notes_text_frame.text.strip()
            if note:
                notes.append((f"Slide {slide_number}", clean_xml_text(note)))
            else:
                notes.append((f"Slide {slide_number}", "No speaker notes."))
        else:
            notes.append((f"Slide {slide_number}", "No speaker notes."))

    return notes


def save_notes_to_docx(notes_list, output_path):
    """Zapisuje notatki do pliku .docx w formie tabeli"""
    doc = Document()
    table = doc.add_table(rows=1, cols=2)
    table.style = 'Table Grid'

    hdr_cells = table.rows[0].cells
    hdr_cells[0].text = 'Slide'
    hdr_cells[1].text = 'Speaker Note'

    for slide_title, note in notes_list:
        row_cells = table.add_row().cells
        row_cells[0].text = slide_title
        row_cells[1].text = note

    doc.save(output_path)
    print(f"✅ Zapisano: {output_path}")


def process_folder(folder_path):
    """Przetwarza wszystkie pliki .pptx i tworzy osobne .docx dla każdego"""
    if not os.path.isdir(folder_path):
        print("❌ Podana ścieżka nie istnieje lub nie jest folderem.")
        return

    pptx_files = [f for f in os.listdir(folder_path) if f.lower().endswith(".pptx")]

    if not pptx_files:
        print("⚠️ Brak plików .pptx w podanym folderze.")
        return

    for filename in pptx_files:
        pptx_path = os.path.join(folder_path, filename)
        notes = extract_notes_from_pptx(pptx_path)

        docx_filename = os.path.splitext(filename)[0] + ".docx"
        docx_path = os.path.join(folder_path, docx_filename)

        save_notes_to_docx(notes, docx_path)


# === 🚀 START SKRYPTU ===
folder_input = input("📁 Podaj ścieżkę do folderu z plikami .pptx: ").strip()
process_folder(folder_input)
