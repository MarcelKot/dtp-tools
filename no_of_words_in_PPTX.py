import os
import re
from pptx import Presentation

def count_words_in_text(text):
    if not text:
        return 0
    return len(re.findall(r'\b\w+\b', text))

def is_slide_hidden(slide):
    return getattr(slide, "hidden", False)

def count_words_in_pptx(file_path, include_notes=False):
    try:
        prs = Presentation(file_path)
        total_words = 0

        # Tekst z widocznych slajdów
        for slide in prs.slides:
            if is_slide_hidden(slide):
                continue
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    total_words += count_words_in_text(shape.text)

        # Notatki prelegenta (opcjonalne)
        if include_notes:
            for slide in prs.slides:
                if is_slide_hidden(slide):
                    continue
                notes_slide = slide.has_notes_slide and slide.notes_slide
                if notes_slide and notes_slide.notes_text_frame:
                    total_words += count_words_in_text(notes_slide.notes_text_frame.text)

        # Tekst z master slide'ów
        for master in prs.slide_masters:
            for shape in master.shapes:
                if hasattr(shape, "text"):
                    total_words += count_words_in_text(shape.text)

        return total_words
    except Exception as e:
        print(f"Błąd podczas przetwarzania pliku {file_path}: {e}")
        return 0

def process_directory(path, include_notes):
    for file_name in os.listdir(path):
        if file_name.lower().endswith(".pptx"):
            full_path = os.path.join(path, file_name)
            word_count = count_words_in_pptx(full_path, include_notes)
            print(f"{file_name}: {word_count} słów")

if __name__ == "__main__":
    folder_path = input("Wklej ścieżkę do folderu z plikami .pptx: ").strip('"').strip()
    notes_input = input("Czy brać pod uwagę speaker notes? (Yes/No): ").strip().lower()
    include_notes = notes_input == "yes"

    if os.path.isdir(folder_path):
        process_directory(folder_path, include_notes)
    else:
        print("Podana ścieżka nie istnieje.")
