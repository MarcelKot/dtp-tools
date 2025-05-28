import os
from pptx import Presentation
from docx import Document


def read_notes_from_docx(docx_path):
    """Czyta notatki z pliku .docx w formacie Slide X | Speaker Note"""
    doc = Document(docx_path)
    notes = {}

    for row in doc.tables[0].rows[1:]:  # Pomijamy nagłówek
        slide_text = row.cells[0].text.strip()
        note = row.cells[1].text.strip()

        if slide_text.lower().startswith("slide"):
            try:
                slide_number = int(slide_text.split()[1])
                if note and note.lower() != "no speaker notes.":
                    notes[slide_number] = note
            except ValueError:
                continue
    return notes


def import_notes_to_pptx(pptx_path, docx_path):
    """Importuje notatki z docx do pptx, nadpisując oryginalny plik"""
    notes_dict = read_notes_from_docx(docx_path)
    prs = Presentation(pptx_path)

    for idx, slide in enumerate(prs.slides):
        slide_number = idx + 1
        if slide_number in notes_dict:
            note = notes_dict[slide_number]
            if not slide.has_notes_slide:
                slide.notes_slide = slide.notes_slide  # wymuszenie utworzenia
            notes_frame = slide.notes_slide.notes_text_frame
            notes_frame.text = note

    prs.save(pptx_path)
    print(f"✅ Nadpisano notatki w: {os.path.basename(pptx_path)}")


def process_folder(folder_path):
    """Skanuje folder i dla każdego pary pptx+docx wykonuje import notatek"""
    if not os.path.isdir(folder_path):
        print("❌ Podana ścieżka nie istnieje lub nie jest folderem.")
        return

    files = os.listdir(folder_path)
    pptx_files = [f for f in files if f.lower().endswith(".pptx")]

    if not pptx_files:
        print("⚠️ Brak plików .pptx w folderze.")
        return

    for pptx_file in pptx_files:
        base_name = os.path.splitext(pptx_file)[0]
        docx_file = base_name + ".docx"

        pptx_path = os.path.join(folder_path, pptx_file)
        docx_path = os.path.join(folder_path, docx_file)

        if os.path.exists(docx_path):
            import_notes_to_pptx(pptx_path, docx_path)
        else:
            print(f"⚠️ Brak pliku .docx dla: {pptx_file} — pomijam.")


# === 🚀 START SKRYPTU ===
folder_input = input("📁 Podaj ścieżkę do folderu z plikami .pptx i .docx: ").strip()
process_folder(folder_input)
